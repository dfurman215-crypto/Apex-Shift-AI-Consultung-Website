from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_picture_shapes(shapes: Iterable):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)


def extract_source_assets(pptx_path: str, output_dir: str | Path) -> tuple[list[dict], str | None]:
    """Extract unique raster images from a source deck and build a labeled contact sheet."""
    source = Path(pptx_path)
    target = Path(output_dir)
    target.mkdir(parents=True, exist_ok=True)

    prs = Presentation(source)
    assets: list[dict] = []
    seen_hashes: set[str] = set()

    for slide_number, slide in enumerate(prs.slides, start=1):
        for picture_index, shape in enumerate(_iter_picture_shapes(slide.shapes), start=1):
            image = shape.image
            blob = image.blob
            digest = hashlib.sha1(blob).hexdigest()
            if digest in seen_hashes:
                continue
            seen_hashes.add(digest)

            ext = (image.ext or "png").lower()
            asset_id = f"A{len(assets) + 1:02d}"
            filename = f"{asset_id}_slide_{slide_number}_{picture_index}.{ext}"
            path = target / filename
            path.write_bytes(blob)

            width = height = None
            try:
                with Image.open(path) as im:
                    width, height = im.size
            except Exception:
                pass

            assets.append(
                {
                    "id": asset_id,
                    "slide": slide_number,
                    "path": str(path),
                    "filename": filename,
                    "width": width,
                    "height": height,
                }
            )

    contact_sheet = _build_contact_sheet(assets, target / "contact_sheet.png") if assets else None
    return assets, str(contact_sheet) if contact_sheet else None


def _build_contact_sheet(assets: list[dict], output_path: Path) -> Path:
    """Create one labeled image so a multimodal local model can inspect source imagery quickly."""
    cols = 3
    cell_w, cell_h = 420, 280
    label_h = 34
    rows = (len(assets) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * cell_w, rows * (cell_h + label_h)), "white")
    draw = ImageDraw.Draw(sheet)

    for index, asset in enumerate(assets):
        row, col = divmod(index, cols)
        x = col * cell_w
        y = row * (cell_h + label_h)
        try:
            with Image.open(asset["path"]) as im:
                tile = ImageOps.contain(im.convert("RGB"), (cell_w - 16, cell_h - 16))
                tx = x + (cell_w - tile.width) // 2
                ty = y + (cell_h - tile.height) // 2
                sheet.paste(tile, (tx, ty))
        except Exception:
            draw.rectangle((x + 8, y + 8, x + cell_w - 8, y + cell_h - 8), outline="gray", width=2)

        label = f"{asset['id']}  |  source slide {asset['slide']}"
        draw.text((x + 12, y + cell_h + 7), label, fill="black")

    sheet.save(output_path, format="PNG", optimize=True)
    return output_path


def asset_map(assets: list[dict] | None) -> dict[str, dict]:
    return {asset["id"]: asset for asset in (assets or [])}
