from pathlib import Path
from typing import Callable

from pptx import Presentation

from agent import LocalAgentUnavailable, plan_deck_with_gemma
from assets import extract_source_assets
from builder import add_slides
from ingest import find_embedded_markdown, inspect_deck, markdown_to_slide_spec


ProgressCallback = Callable[[str], None] | None
IMAGE_FRIENDLY_TYPES = {"hero", "metaphor", "ecosystem", "statement"}


def _assign_source_assets(spec: dict, assets: list[dict]) -> None:
    """Reuse a few source-deck images deterministically without asking Gemma to inspect them."""
    if not assets:
        return

    preferred = sorted(assets, key=lambda asset: (asset.get("slide", 999), asset.get("id", "")))
    asset_ids = [asset["id"] for asset in preferred]
    used = set()
    cursor = 0

    for slide in spec.get("slides", []):
        if slide.get("type") not in IMAGE_FRIENDLY_TYPES:
            continue
        if slide.get("asset_id"):
            continue

        while cursor < len(asset_ids) and asset_ids[cursor] in used:
            cursor += 1
        if cursor >= len(asset_ids):
            break

        slide["asset_id"] = asset_ids[cursor]
        used.add(asset_ids[cursor])
        cursor += 1


def amend_deck(
    source_pptx: str,
    output_pptx: str,
    use_agent: bool = True,
    progress_callback: ProgressCallback = None,
) -> dict:
    source = Path(source_pptx)
    output = Path(output_pptx)

    def report(message: str):
        if progress_callback:
            progress_callback(message)

    report("Validating source and output files...")
    if not source.exists():
        raise FileNotFoundError(source)
    if source.resolve() == output.resolve():
        raise ValueError("Output must be a different file so the source deck is not overwritten.")

    report("Locating the embedded Markdown presentation brief...")
    markdown_slide, markdown_text = find_embedded_markdown(str(source))

    report("Extracting reusable source-deck images...")
    asset_dir = output.parent / ".present_assets" / source.stem
    assets, _contact_sheet = extract_source_assets(str(source), asset_dir)

    planner = "deterministic parser"
    agent_error = None

    if use_agent:
        try:
            report("Summarizing the existing deck for the local Gemma planner...")
            deck_info = inspect_deck(str(source))

            report("Asking local Gemma for a fast text-only slide plan...")
            spec = plan_deck_with_gemma(deck_info, markdown_text)
            planner = "local Gemma via Ollama (fast text-only plan)"
        except (LocalAgentUnavailable, ValueError) as exc:
            agent_error = str(exc)
            report("Gemma planning was unavailable; using PRESENT's deterministic fallback...")
            spec = markdown_to_slide_spec(markdown_text)
    else:
        report("Parsing the embedded Markdown presentation blueprint...")
        spec = markdown_to_slide_spec(markdown_text)

    report("Assigning reusable source imagery to selected visual slides...")
    _assign_source_assets(spec, assets)

    report("Opening the source PowerPoint and preserving existing slides...")
    prs = Presentation(source)
    original_slide_count = len(prs.slides)

    report("Composing the new visual slides...")
    add_slides(prs, spec, assets=assets)

    report("Saving the amended PowerPoint...")
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    report("Build complete.")
    return {
        "source": str(source),
        "output": str(output),
        "markdown_slide": markdown_slide,
        "original_slide_count": original_slide_count,
        "added_slide_count": len(spec.get("slides", [])),
        "final_slide_count": len(prs.slides),
        "planner": planner,
        "agent_error": agent_error,
        "asset_count": len(assets),
    }
