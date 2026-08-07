import re
from pathlib import Path
from pptx import Presentation


def _normalize_powerpoint_text(text: str) -> str:
    """Normalize PowerPoint paragraph/soft-break characters for parsing."""
    return (
        text.replace("\x0b", "\n")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
    )


def extract_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = _normalize_powerpoint_text(shape.text).strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def inspect_deck(pptx_path: str) -> dict:
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(path)

    prs = Presentation(path)
    slides = []
    for number, slide in enumerate(prs.slides, start=1):
        slides.append({
            "number": number,
            "text": extract_slide_text(slide),
            "shape_count": len(slide.shapes),
        })

    return {
        "path": str(path),
        "slide_count": len(prs.slides),
        "slides": slides,
    }


def find_embedded_markdown(pptx_path: str) -> tuple[int, str]:
    prs = Presentation(pptx_path)
    candidates = []

    for number, slide in enumerate(prs.slides, start=1):
        text = extract_slide_text(slide)
        score = 0
        if "# Apex Shift PEARL Playbook" in text:
            score += 100
        if "# Slide Blueprint" in text:
            score += 50
        if "## Purpose" in text:
            score += 25
        if "## " in text:
            score += 10
        if len(text) > 1000:
            score += 5
        if score:
            candidates.append((score, number, text))

    if not candidates:
        raise ValueError("No embedded Markdown presentation brief was found in the deck.")

    candidates.sort(reverse=True)
    _, number, text = candidates[0]
    return number, text


def _clean_lines(block: str) -> list[str]:
    cleaned = []
    for raw in block.splitlines():
        line = raw.strip()
        if not line:
            continue
        if set(line) == {"-"}:
            continue
        line = re.sub(r"^[-*]\s*", "", line)
        cleaned.append(line)
    return cleaned


def markdown_to_slide_spec(markdown_text: str) -> dict:
    markdown_text = _normalize_powerpoint_text(markdown_text)

    blueprint_match = re.search(
        r"# Slide Blueprint\s*(.*?)(?:# Presentation Style Guide|\Z)",
        markdown_text,
        flags=re.S | re.I,
    )
    blueprint = blueprint_match.group(1) if blueprint_match else markdown_text

    sections = re.split(r"(?m)^##\s+(?=\d+\.)", blueprint)
    slides = []

    for section in sections:
        section = section.strip()
        if not section:
            continue

        first_line, *rest = section.splitlines()
        heading_match = re.match(r"(\d+)\.\s*(.+)", first_line.strip())
        if not heading_match:
            continue

        number = int(heading_match.group(1))
        section_name = heading_match.group(2).strip()
        lines = _clean_lines("\n".join(rest))

        fields = {}
        body_lines = []
        for line in lines:
            key_match = re.match(
                r"^(Title|Subtitle|Headline|Subheadline|Message|Visual|Content|Center|Explain):\s*(.*)$",
                line,
                re.I,
            )
            if key_match:
                fields[key_match.group(1).lower()] = key_match.group(2).strip().strip("*")
            else:
                body_lines.append(line.strip("*"))

        headline = fields.get("title") or fields.get("headline") or section_name
        subheadline = fields.get("subtitle") or fields.get("subheadline")

        joined = " ".join(lines).lower()
        if number == 1 or section_name.lower() == "cover":
            slide = {
                "type": "hero",
                "headline": headline,
                "subheadline": subheadline or "Pilot • Engagement • Roadmap • Layering",
            }
        elif "traditional" in joined and "pearl" in joined:
            slide = {
                "type": "comparison",
                "headline": headline,
                "left_title": "Traditional Project",
                "left": ["Large scope", "High risk", "Long delivery", "Low visibility"],
                "right_title": "PEARL",
                "right": ["Small pilot", "Progressive value", "Continuous collaboration", "Reduced risk"],
            }
        elif number == 13 or section_name.lower() == "closing":
            slide = {
                "type": "statement",
                "headline": headline,
                "subheadline": subheadline or "Every business transformation begins with a single opportunity.",
            }
        else:
            items = []
            for value in fields.values():
                if value and value not in (headline, subheadline):
                    items.append(value)
            items.extend(body_lines)
            slide = {
                "type": "content",
                "headline": headline,
                "subheadline": subheadline,
                "items": items[:12],
            }

        slide["source_section"] = number
        slide["source_name"] = section_name
        slides.append(slide)

    if not slides:
        raise ValueError("The embedded Markdown did not contain recognizable numbered slide blueprint sections.")

    return {
        "deck": "Embedded Markdown Build",
        "version": "0.1",
        "slides": slides,
    }
