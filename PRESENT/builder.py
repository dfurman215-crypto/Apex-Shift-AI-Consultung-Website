from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import design
from assets import asset_map


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        try:
            if len(layout.placeholders) == 0:
                return layout
        except Exception:
            pass
    return min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))


def _new_slide(prs, background=None):
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background or design.PEARL_WHITE
    return slide


def _style_paragraph(paragraph, size, color, bold=False, font_name=None, align=None):
    paragraph.font.size = size
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = font_name or design.FONT_BODY
    if align is not None:
        paragraph.alignment = align


def _textbox(slide, text, x, y, w, h, *, size=None, color=None, bold=False,
             font_name=None, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, rotation=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.rotation = rotation
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = str(text or "")
    _style_paragraph(
        p,
        size or design.BODY_SIZE,
        color or design.CHARCOAL,
        bold=bold,
        font_name=font_name,
        align=align,
    )
    return box


def _shape(slide, shape_type, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def _card(slide, x, y, w, h, *, fill=design.WHITE, line=design.LIGHT_GRAY):
    return _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def _title(slide, text, *, color=design.CHARCOAL, kicker=None):
    if kicker:
        _textbox(slide, kicker.upper(), 0.75, 0.36, 4.8, 0.28,
                 size=design.MICRO_SIZE, color=design.GOLD, bold=True)
    _textbox(slide, text, 0.72, 0.58, 11.8, 0.74,
             size=design.TITLE_SIZE, color=color, bold=True, font_name=design.FONT_HEADLINE)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.73, 1.34, 0.72, 0.055, design.GOLD)


def _subtitle(slide, text, y=1.44, w=11.5, color=design.MUTED):
    if text:
        _textbox(slide, text, 0.74, y, w, 0.52, size=design.SUBTITLE_SIZE, color=color)


def _bullet_list(slide, items, x, y, w, h, *, color=design.CHARCOAL,
                 dot_color=design.SOFT_TEAL, size=design.BODY_SIZE, max_items=6):
    items = [str(item) for item in (items or []) if str(item).strip()][:max_items]
    if not items:
        return
    row_h = h / max(len(items), 1)
    for idx, item in enumerate(items):
        cy = y + idx * row_h + row_h * 0.22
        _shape(slide, MSO_SHAPE.OVAL, x, cy, 0.10, 0.10, dot_color)
        _textbox(slide, item, x + 0.2, y + idx * row_h, w - 0.2, row_h,
                 size=size, color=color, valign=MSO_ANCHOR.MIDDLE)


def _picture_fill(slide, path, x, y, w, h):
    path = Path(path)
    if not path.exists():
        return None
    try:
        with Image.open(path) as im:
            im = im.convert("RGB")
            target_ratio = max(w / h, 0.01)
            target_w = 1600
            target_h = max(1, int(target_w / target_ratio))
            fitted = ImageOps.fit(im, (target_w, target_h), method=Image.Resampling.LANCZOS)
            stream = BytesIO()
            fitted.save(stream, format="JPEG", quality=92)
            stream.seek(0)
            return slide.shapes.add_picture(stream, Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def _asset_path(slide_spec, assets):
    if not assets:
        return None
    lookup = asset_map(assets)
    asset_id = slide_spec.get("asset_id")
    if asset_id and asset_id in lookup:
        return lookup[asset_id]["path"]
    return None


def _pill(slide, text, x, y, w, *, fill=design.PALE_BLUE, color=design.DEEP_BLUE):
    _card(slide, x, y, w, 0.42, fill=fill, line=None)
    _textbox(slide, text, x + 0.05, y + 0.02, w - 0.1, 0.36,
             size=design.SMALL_SIZE, color=color, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def build_hero(prs, slide_spec, assets=None):
    path = _asset_path(slide_spec, assets)
    slide = _new_slide(prs, design.DEEP_BLUE if path else design.PEARL_WHITE)

    if path:
        _picture_fill(slide, path, 7.35, 0, 5.98, 7.5)
        _shape(slide, MSO_SHAPE.RECTANGLE, 7.25, 0, 0.12, 7.5, design.GOLD)
        _textbox(slide, "APEX SHIFT", 0.78, 0.7, 3.1, 0.32,
                 size=design.SMALL_SIZE, color=design.GOLD, bold=True)
        _textbox(slide, slide_spec.get("headline", ""), 0.78, 1.35, 5.95, 2.1,
                 size=design.HERO_SIZE, color=design.WHITE, bold=True,
                 font_name=design.FONT_HEADLINE, valign=MSO_ANCHOR.MIDDLE)
        _textbox(slide, slide_spec.get("subheadline", ""), 0.82, 3.65, 5.8, 1.2,
                 size=design.SUBTITLE_SIZE, color=design.PEARL_WHITE)
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 5.35, 1.0, 0.07, design.GOLD)
    else:
        _textbox(slide, "APEX SHIFT", 0.82, 0.72, 3.0, 0.32,
                 size=design.SMALL_SIZE, color=design.GOLD, bold=True)
        _textbox(slide, slide_spec.get("headline", ""), 0.82, 1.55, 10.9, 1.7,
                 size=design.HERO_SIZE, color=design.CHARCOAL, bold=True,
                 font_name=design.FONT_HEADLINE)
        _textbox(slide, slide_spec.get("subheadline", ""), 0.86, 3.35, 9.8, 1.0,
                 size=design.SUBTITLE_SIZE, color=design.MUTED)
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.86, 4.7, 1.1, 0.07, design.GOLD)
    return slide


def build_metaphor(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Concept")
    _subtitle(slide, slide_spec.get("subheadline"))
    path = _asset_path(slide_spec, assets)

    if path:
        _picture_fill(slide, path, 0.72, 2.02, 6.05, 4.82)
        _card(slide, 7.12, 2.03, 5.48, 4.82, fill=design.WHITE)
        message = slide_spec.get("message") or slide_spec.get("body") or slide_spec.get("items", [])
        if isinstance(message, list):
            _bullet_list(slide, message, 7.48, 2.58, 4.72, 3.45, max_items=5)
        else:
            _textbox(slide, message, 7.5, 2.55, 4.62, 2.7,
                     size=Pt(22), color=design.DEEP_BLUE, bold=True,
                     font_name=design.FONT_HEADLINE, valign=MSO_ANCHOR.MIDDLE)
    else:
        build_content_into(slide, slide_spec, top=2.0)
    return slide


def build_comparison(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Comparison")
    _subtitle(slide, slide_spec.get("subheadline"))

    _card(slide, 0.74, 2.03, 5.7, 4.75, fill=design.PALE_BLUE)
    _card(slide, 6.89, 2.03, 5.7, 4.75, fill=design.LIGHT_TEAL)
    _shape(slide, MSO_SHAPE.OVAL, 6.18, 3.72, 0.96, 0.96, design.DEEP_BLUE)
    _textbox(slide, "VS", 6.18, 3.73, 0.96, 0.9, size=design.SMALL_SIZE,
             color=design.WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _textbox(slide, slide_spec.get("left_title", "Current"), 1.08, 2.4, 4.8, 0.55,
             size=Pt(23), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE)
    _bullet_list(slide, slide_spec.get("left", []), 1.08, 3.12, 4.85, 2.95,
                 dot_color=design.MID_BLUE, max_items=5)

    _textbox(slide, slide_spec.get("right_title", "Target"), 7.24, 2.4, 4.8, 0.55,
             size=Pt(23), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE)
    _bullet_list(slide, slide_spec.get("right", []), 7.24, 3.12, 4.85, 2.95,
                 dot_color=design.SOFT_TEAL, max_items=5)
    return slide


def build_content_into(slide, slide_spec, top=1.95):
    items = [str(x) for x in slide_spec.get("items", []) if str(x).strip()][:6]
    if not items:
        items = [str(slide_spec.get("body", ""))]
    count = len(items)
    cols = 3 if count >= 5 else 2
    rows = (count + cols - 1) // cols
    gap_x, gap_y = 0.24, 0.24
    total_w = 11.86
    card_w = (total_w - gap_x * (cols - 1)) / cols
    avail_h = 6.72 - top
    card_h = min(1.75, (avail_h - gap_y * (rows - 1)) / rows)

    for idx, item in enumerate(items):
        row, col = divmod(idx, cols)
        x = 0.74 + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)
        _card(slide, x, y, card_w, card_h, fill=design.WHITE)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, card_h, design.SOFT_TEAL)
        _textbox(slide, f"{idx + 1:02d}", x + 0.24, y + 0.23, 0.55, 0.3,
                 size=design.MICRO_SIZE, color=design.GOLD, bold=True)
        _textbox(slide, item, x + 0.24, y + 0.55, card_w - 0.48, card_h - 0.72,
                 size=design.CARD_TITLE_SIZE, color=design.CHARCOAL, bold=True,
                 font_name=design.FONT_HEADLINE, valign=MSO_ANCHOR.MIDDLE)


def build_content(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Executive View")
    _subtitle(slide, slide_spec.get("subheadline"))
    build_content_into(slide, slide_spec, top=2.0 if slide_spec.get("subheadline") else 1.72)
    return slide


def _stage_label(stage):
    if isinstance(stage, dict):
        return str(stage.get("label") or stage.get("title") or stage.get("name") or "")
    return str(stage)


def _stage_detail(stage):
    if isinstance(stage, dict):
        return str(stage.get("detail") or stage.get("description") or "")
    return ""


def build_process(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Method")
    _subtitle(slide, slide_spec.get("subheadline"))
    stages = list(slide_spec.get("stages") or slide_spec.get("items") or [])[:5]
    if not stages:
        build_content_into(slide, slide_spec, top=2.0)
        return slide

    n = len(stages)
    x0, y = 0.82, 2.5
    usable_w = 11.7
    gap = 0.28
    card_w = (usable_w - gap * (n - 1)) / n

    for idx, stage in enumerate(stages):
        x = x0 + idx * (card_w + gap)
        if idx < n - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + card_w), Inches(y + 1.35),
                Inches(x + card_w + gap), Inches(y + 1.35),
            )
            conn.line.color.rgb = design.GOLD
            conn.line.width = Pt(2)
        _card(slide, x, y, card_w, 2.8, fill=design.WHITE)
        _shape(slide, MSO_SHAPE.OVAL, x + card_w / 2 - 0.34, y - 0.34, 0.68, 0.68,
               design.DEEP_BLUE)
        _textbox(slide, str(idx + 1), x + card_w / 2 - 0.34, y - 0.31, 0.68, 0.58,
                 size=design.SMALL_SIZE, color=design.WHITE, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _textbox(slide, _stage_label(stage), x + 0.16, y + 0.5, card_w - 0.32, 0.72,
                 size=design.CARD_TITLE_SIZE, color=design.DEEP_BLUE, bold=True,
                 font_name=design.FONT_HEADLINE, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        detail = _stage_detail(stage)
        if detail:
            _textbox(slide, detail, x + 0.18, y + 1.38, card_w - 0.36, 1.05,
                     size=Pt(12), color=design.MUTED, align=PP_ALIGN.CENTER)
    return slide


def build_lifecycle(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Capability Lifecycle")
    _subtitle(slide, slide_spec.get("subheadline"))
    stages = list(slide_spec.get("stages") or slide_spec.get("items") or [])[:5]
    if not stages:
        build_content_into(slide, slide_spec, top=2.0)
        return slide

    n = len(stages)
    centers = []
    x0, x1 = 1.45, 11.9
    for idx in range(n):
        x = x0 + (x1 - x0) * idx / max(n - 1, 1)
        centers.append(x)

    baseline_y = 3.75
    for idx in range(n - 1):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(centers[idx] + 0.55), Inches(baseline_y),
            Inches(centers[idx + 1] - 0.55), Inches(baseline_y),
        )
        conn.line.color.rgb = design.GOLD
        conn.line.width = Pt(2.5)

    sizes = [0.86 + 0.18 * idx for idx in range(n)]
    fills = [design.WHITE, design.PALE_BLUE, design.LIGHT_TEAL, design.LIGHT_GOLD, design.WHITE]
    for idx, stage in enumerate(stages):
        size = sizes[idx]
        x = centers[idx] - size / 2
        y = baseline_y - size / 2
        pearl = _shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fills[idx % len(fills)], design.GOLD)
        pearl.line.width = Pt(1.4)
        _textbox(slide, str(idx + 1), centers[idx] - 0.22, baseline_y - 0.17, 0.44, 0.34,
                 size=design.SMALL_SIZE, color=design.DEEP_BLUE, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _textbox(slide, _stage_label(stage), centers[idx] - 1.15, 4.72, 2.3, 0.62,
                 size=design.CARD_TITLE_SIZE, color=design.DEEP_BLUE, bold=True,
                 font_name=design.FONT_HEADLINE, align=PP_ALIGN.CENTER)
        detail = _stage_detail(stage)
        if detail:
            _textbox(slide, detail, centers[idx] - 1.1, 5.42, 2.2, 0.75,
                     size=Pt(11), color=design.MUTED, align=PP_ALIGN.CENTER)
    return slide


def build_collaboration(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Co-Creation")
    _subtitle(slide, slide_spec.get("subheadline"))

    left_title = slide_spec.get("left_title", "Apex Shift")
    right_title = slide_spec.get("right_title", "Client")
    center = slide_spec.get("center", "Collaboration")

    _card(slide, 0.8, 2.35, 3.5, 3.65, fill=design.PALE_BLUE)
    _card(slide, 9.03, 2.35, 3.5, 3.65, fill=design.LIGHT_TEAL)
    _shape(slide, MSO_SHAPE.OVAL, 5.25, 3.0, 2.82, 2.82, design.DEEP_BLUE)

    for x1, x2 in ((4.3, 5.25), (8.07, 9.03)):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(4.41), Inches(x2), Inches(4.41)
        )
        conn.line.color.rgb = design.GOLD
        conn.line.width = Pt(2.5)

    _textbox(slide, left_title, 1.15, 2.7, 2.8, 0.55,
             size=Pt(22), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE,
             align=PP_ALIGN.CENTER)
    _bullet_list(slide, slide_spec.get("left_items", slide_spec.get("left", [])),
                 1.15, 3.45, 2.72, 2.0, max_items=4, size=Pt(13))

    _textbox(slide, center, 5.55, 3.68, 2.22, 0.9,
             size=Pt(20), color=design.WHITE, bold=True, font_name=design.FONT_HEADLINE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _textbox(slide, slide_spec.get("center_detail", "Joint innovation"), 5.56, 4.58, 2.2, 0.55,
             size=Pt(11), color=design.PEARL_WHITE, align=PP_ALIGN.CENTER)

    _textbox(slide, right_title, 9.38, 2.7, 2.8, 0.55,
             size=Pt(22), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE,
             align=PP_ALIGN.CENTER)
    _bullet_list(slide, slide_spec.get("right_items", slide_spec.get("right", [])),
                 9.38, 3.45, 2.72, 2.0, max_items=4, size=Pt(13))
    return slide


def _roadmap_lane(slide, title, items, y, accent):
    _textbox(slide, title, 0.84, y - 0.18, 2.0, 0.5,
             size=design.CARD_TITLE_SIZE, color=accent, bold=True, font_name=design.FONT_HEADLINE)
    start_x, end_x = 3.0, 12.0
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(start_x), Inches(y + 0.25), Inches(end_x), Inches(y + 0.25)
    )
    conn.line.color.rgb = accent
    conn.line.width = Pt(2.2)
    items = list(items or [])[:5]
    for idx, item in enumerate(items):
        x = start_x + (end_x - start_x) * idx / max(len(items) - 1, 1)
        _shape(slide, MSO_SHAPE.OVAL, x - 0.17, y + 0.08, 0.34, 0.34, accent)
        _textbox(slide, _stage_label(item), x - 0.9, y + 0.55, 1.8, 0.7,
                 size=Pt(11), color=design.CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_roadmap(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Roadmap Alignment")
    _subtitle(slide, slide_spec.get("subheadline"))
    _roadmap_lane(slide, slide_spec.get("strategic_title", "Strategic Roadmap"),
                  slide_spec.get("strategic", []), 2.75, design.DEEP_BLUE)
    _roadmap_lane(slide, slide_spec.get("tactical_title", "Tactical Roadmap"),
                  slide_spec.get("tactical", []), 4.75, design.SOFT_TEAL)
    cadence = slide_spec.get("cadence")
    if cadence:
        _pill(slide, cadence, 4.45, 6.42, 4.45, fill=design.LIGHT_GOLD, color=design.CHARCOAL)
    return slide


def build_layers(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Layering")
    _subtitle(slide, slide_spec.get("subheadline"))

    center = slide_spec.get("center", "Pilot")
    layers = list(slide_spec.get("layers") or slide_spec.get("items") or [])[:4]
    diameters = [4.3, 3.45, 2.6, 1.75]
    fills = [design.PALE_BLUE, design.LIGHT_TEAL, design.LIGHT_GOLD, design.WHITE]
    cx, cy = 3.55, 4.4
    for idx, diameter in enumerate(diameters[: max(1, len(layers) + 1)]):
        _shape(slide, MSO_SHAPE.OVAL, cx - diameter / 2, cy - diameter / 2,
               diameter, diameter, fills[idx % len(fills)], design.GOLD)
    _shape(slide, MSO_SHAPE.OVAL, cx - 0.58, cy - 0.58, 1.16, 1.16, design.DEEP_BLUE)
    _textbox(slide, center, cx - 0.46, cy - 0.28, 0.92, 0.56,
             size=Pt(13), color=design.WHITE, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)

    _textbox(slide, "Progressive value", 6.55, 2.25, 4.9, 0.5,
             size=Pt(22), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE)
    for idx, layer in enumerate(layers):
        y = 3.0 + idx * 0.83
        _shape(slide, MSO_SHAPE.OVAL, 6.62, y + 0.07, 0.28, 0.28,
               [design.GOLD, design.SOFT_TEAL, design.MID_BLUE, design.DEEP_BLUE][idx % 4])
        _textbox(slide, _stage_label(layer), 7.12, y - 0.03, 4.65, 0.42,
                 size=design.CARD_TITLE_SIZE, color=design.CHARCOAL, bold=True,
                 font_name=design.FONT_HEADLINE)
        detail = _stage_detail(layer)
        if detail:
            _textbox(slide, detail, 7.12, y + 0.35, 4.65, 0.38,
                     size=Pt(11), color=design.MUTED)
    return slide


def build_ecosystem(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Integration")
    _subtitle(slide, slide_spec.get("subheadline"))
    path = _asset_path(slide_spec, assets)

    center_x, center_y = (4.3 if path else 6.65), 4.25
    nodes = list(slide_spec.get("nodes") or slide_spec.get("items") or [])[:6]
    if not nodes:
        nodes = ["Capability", "Data", "Workflow", "People", "Systems", "Next Pearl"]

    positions = [
        (center_x - 2.65, center_y - 1.45),
        (center_x, center_y - 1.9),
        (center_x + 2.65, center_y - 1.45),
        (center_x - 2.65, center_y + 1.35),
        (center_x, center_y + 1.75),
        (center_x + 2.65, center_y + 1.35),
    ]
    for x, y in positions[: len(nodes)]:
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(center_x), Inches(center_y), Inches(x), Inches(y)
        )
        conn.line.color.rgb = design.LIGHT_GRAY
        conn.line.width = Pt(1.5)

    _shape(slide, MSO_SHAPE.OVAL, center_x - 0.83, center_y - 0.83, 1.66, 1.66, design.DEEP_BLUE)
    _textbox(slide, slide_spec.get("center", "Integrated\nSystem"), center_x - 0.68, center_y - 0.45,
             1.36, 0.9, size=Pt(14), color=design.WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    for idx, node in enumerate(nodes):
        x, y = positions[idx]
        _shape(slide, MSO_SHAPE.OVAL, x - 0.68, y - 0.42, 1.36, 0.84,
               design.WHITE, design.SOFT_TEAL)
        _textbox(slide, _stage_label(node), x - 0.57, y - 0.28, 1.14, 0.56,
                 size=Pt(10), color=design.CHARCOAL, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

    if path:
        _picture_fill(slide, path, 9.2, 2.15, 3.4, 4.9)
    return slide


def build_pricing(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Pricing Philosophy")
    _subtitle(slide, slide_spec.get("subheadline"))

    formula = slide_spec.get("formula") or ["Scope", "Effort", "Price"]
    if isinstance(formula, str):
        formula = [part.strip() for part in formula.replace("→", ">").split(">") if part.strip()]
    formula = list(formula)[:4]
    x0, y, box_w, gap = 1.2, 2.15, 2.55, 0.55
    for idx, item in enumerate(formula):
        x = x0 + idx * (box_w + gap)
        _card(slide, x, y, box_w, 1.2, fill=design.DEEP_BLUE if idx == len(formula) - 1 else design.WHITE,
              line=design.DEEP_BLUE)
        _textbox(slide, item, x + 0.15, y + 0.22, box_w - 0.3, 0.72,
                 size=Pt(21), color=design.WHITE if idx == len(formula) - 1 else design.DEEP_BLUE,
                 bold=True, font_name=design.FONT_HEADLINE, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        if idx < len(formula) - 1:
            _textbox(slide, "→", x + box_w + 0.08, y + 0.28, 0.4, 0.55,
                     size=Pt(24), color=design.GOLD, bold=True, align=PP_ALIGN.CENTER)

    factors = [str(x) for x in slide_spec.get("factors", slide_spec.get("items", []))][:6]
    if factors:
        _textbox(slide, "What drives effort", 0.86, 3.85, 3.0, 0.42,
                 size=design.CARD_TITLE_SIZE, color=design.MUTED, bold=True)
        for idx, factor in enumerate(factors):
            col, row = idx % 3, idx // 3
            _pill(slide, factor, 0.86 + col * 4.05, 4.42 + row * 0.72, 3.65,
                  fill=design.PALE_BLUE if row == 0 else design.LIGHT_TEAL,
                  color=design.DEEP_BLUE)
    payment = slide_spec.get("payment") or slide_spec.get("footer")
    if payment:
        _card(slide, 0.86, 6.2, 11.6, 0.62, fill=design.LIGHT_GOLD, line=None)
        _textbox(slide, payment, 1.02, 6.29, 11.28, 0.42,
                 size=Pt(12), color=design.CHARCOAL, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
    return slide


def build_benchmark(prs, slide_spec, assets=None):
    slide = _new_slide(prs)
    _title(slide, slide_spec.get("headline", ""), kicker="Market Context")
    _subtitle(slide, slide_spec.get("subheadline"))
    metrics = slide_spec.get("metrics") or []
    if not metrics:
        metrics = [{"label": f"Benchmark {i+1}", "value": item} for i, item in enumerate(slide_spec.get("items", [])[:4])]
    metrics = list(metrics)[:4]
    if not metrics:
        build_content_into(slide, slide_spec, top=2.0)
        return slide

    n = len(metrics)
    gap = 0.28
    card_w = (11.75 - gap * (n - 1)) / n
    for idx, metric in enumerate(metrics):
        x = 0.79 + idx * (card_w + gap)
        _card(slide, x, 2.2, card_w, 3.85, fill=design.WHITE)
        label = metric.get("label", "Benchmark") if isinstance(metric, dict) else "Benchmark"
        value = metric.get("value", "") if isinstance(metric, dict) else str(metric)
        detail = metric.get("detail", "") if isinstance(metric, dict) else ""
        _textbox(slide, label.upper(), x + 0.18, 2.55, card_w - 0.36, 0.42,
                 size=design.MICRO_SIZE, color=design.GOLD, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, value, x + 0.18, 3.18, card_w - 0.36, 1.0,
                 size=Pt(24), color=design.DEEP_BLUE, bold=True, font_name=design.FONT_HEADLINE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if detail:
            _textbox(slide, detail, x + 0.22, 4.5, card_w - 0.44, 0.95,
                     size=Pt(11), color=design.MUTED, align=PP_ALIGN.CENTER)
    return slide


def build_statement(prs, slide_spec, assets=None):
    path = _asset_path(slide_spec, assets)
    slide = _new_slide(prs, design.DEEP_BLUE)
    if path:
        _picture_fill(slide, path, 8.55, 0, 4.78, 7.5)
        _shape(slide, MSO_SHAPE.RECTANGLE, 8.43, 0, 0.12, 7.5, design.GOLD)
        text_w = 6.9
    else:
        text_w = 11.0
    _textbox(slide, slide_spec.get("headline", ""), 0.98, 2.15, text_w, 1.65,
             size=design.HERO_SIZE, color=design.WHITE, bold=True,
             font_name=design.FONT_HEADLINE, valign=MSO_ANCHOR.MIDDLE)
    if slide_spec.get("subheadline"):
        _textbox(slide, slide_spec.get("subheadline"), 1.03, 4.02, min(text_w, 8.7), 1.0,
                 size=design.SUBTITLE_SIZE, color=design.PEARL_WHITE)
    _shape(slide, MSO_SHAPE.RECTANGLE, 1.03, 5.42, 1.0, 0.07, design.GOLD)
    return slide


BUILDERS = {
    "hero": build_hero,
    "metaphor": build_metaphor,
    "comparison": build_comparison,
    "content": build_content,
    "process": build_process,
    "lifecycle": build_lifecycle,
    "collaboration": build_collaboration,
    "roadmap": build_roadmap,
    "layers": build_layers,
    "ecosystem": build_ecosystem,
    "pricing": build_pricing,
    "benchmark": build_benchmark,
    "statement": build_statement,
}


def _apply_asset_defaults(slides, assets):
    if not assets:
        return
    preferred = [asset["id"] for asset in assets[:6]]
    cursor = 0
    visual_types = {"hero", "metaphor", "ecosystem", "statement"}
    for slide in slides:
        if slide.get("asset_id"):
            continue
        if slide.get("type") in visual_types and cursor < len(preferred):
            slide["asset_id"] = preferred[cursor]
            cursor += 1


def add_slides(prs, spec: dict, assets=None):
    slides = list(spec.get("slides", []))
    _apply_asset_defaults(slides, assets)
    for slide_spec in slides:
        slide_type = slide_spec.get("type", "content")
        builder = BUILDERS.get(slide_type, build_content)
        builder(prs, slide_spec, assets)
    return prs


def build_deck(spec: dict, output_path: str, assets=None):
    prs = Presentation()
    prs.slide_width = design.SLIDE_WIDTH
    prs.slide_height = design.SLIDE_HEIGHT
    add_slides(prs, spec, assets=assets)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
